# -*- coding: utf-8 -*-
"""Boss-battle deck: 「なぜ私たち」＋ ひろゆき(論破王FM)ボス戦.
Same dark MBB design tokens as build_deck.py. Format: Title / KeyMessage / Body.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
BG_COVER = os.path.join(HERE, "assets", "image1.png")
BG_BODY  = os.path.join(HERE, "assets", "image13.png")
OUT = os.path.join(os.path.dirname(HERE), "outputs", "07_ひろゆきボス戦.pptx")

CY="5BC8FF"; WH="EAF1FF"; MUT="9DB0D0"; DIM="6E7FA6"; GOLD="F5B454"; PUR="9C8CFF"
CARD="111C3C"; CARDLN="2A3B66"; KMFILL="13234D"; RED="FF5C6C"; REDBG="2A1530"
FONT="Arial"; EAFONT="Meiryo"

prs=Presentation(); prs.slide_width=Inches(10); prs.slide_height=Inches(5.625)
BLANK=prs.slide_layouts[6]
def rgb(h): return RGBColor.from_string(h)

def slide(bg=BG_BODY):
    s=prs.slides.add_slide(BLANK)
    pic=s.shapes.add_picture(bg,0,0,prs.slide_width,prs.slide_height)
    sp=pic._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2,sp)
    return s

def _font(run,size,color,bold,font=FONT):
    f=run.font; f.size=Pt(size); f.bold=bold; f.name=font; f.color.rgb=rgb(color)
    rPr=run._r.get_or_add_rPr(); ea=rPr.find(qn('a:ea'))
    if ea is None: ea=rPr.makeelement(qn('a:ea'),{}); rPr.append(ea)
    ea.set('typeface',EAFONT)

def text(s,x,y,w,h,runs,size=12,color=MUT,bold=False,align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP,ls=1.0,font=FONT):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    tf.vertical_anchor=anchor
    if isinstance(runs,str): paras=[[(runs,color,bold,size)]]
    elif runs and isinstance(runs[0],tuple): paras=[runs]
    else: paras=runs
    for pi,para in enumerate(paras):
        p=tf.paragraphs[0] if pi==0 else tf.add_paragraph()
        p.alignment=align; p.line_spacing=ls
        for tup in para:
            txt=tup[0]; c=tup[1] if len(tup)>1 and tup[1] else color
            b=tup[2] if len(tup)>2 and tup[2] is not None else bold
            sz=tup[3] if len(tup)>3 and tup[3] else size
            r=p.add_run(); r.text=txt; _font(r,sz,c,b,font)
    return tb

def _shadow(shape):
    spPr=shape._element.spPr; el=spPr.makeelement(qn('a:effectLst'),{})
    sh=spPr.makeelement(qn('a:outerShdw'),{'blurRad':'90000','dist':'30000','dir':'5400000','rotWithShape':'0','algn':'tl'})
    clr=spPr.makeelement(qn('a:srgbClr'),{'val':'000000'}); al=spPr.makeelement(qn('a:alpha'),{'val':'38000'})
    clr.append(al); sh.append(clr); el.append(sh); spPr.append(el)

def card(s,x,y,w,h,fill=CARD,line=CARDLN,radius=0.07,shadow=True,line_w=1.0):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=rgb(fill)
    if line: sh.line.color.rgb=rgb(line); sh.line.width=Pt(line_w)
    else: sh.line.fill.background()
    try: sh.adjustments[0]=radius
    except: pass
    sh.shadow.inherit=False
    if shadow:_shadow(sh)
    if sh.has_text_frame: sh.text_frame.word_wrap=True
    return sh

def rect(s,x,y,w,h,color):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=rgb(color); sh.line.fill.background(); sh.shadow.inherit=False
    return sh

def oval(s,x,y,w,h,fill,line=None,line_w=2.0):
    sh=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=rgb(fill)
    if line: sh.line.color.rgb=rgb(line); sh.line.width=Pt(line_w)
    else: sh.line.fill.background()
    sh.shadow.inherit=False; return sh

def header(s,kicker,title,ksize=27,kc=CY):
    text(s,0.6,0.34,8.8,0.28,kicker,size=11,color=kc,bold=True)
    text(s,0.6,0.62,8.8,0.66,title,size=ksize,color=WH,bold=True,ls=1.0)

def keymsg(s,runs,y=1.34,h=0.60,barc=GOLD):
    card(s,0.6,y,8.8,h,fill=KMFILL,line="2C4A86",radius=0.10,shadow=False,line_w=1.0)
    rect(s,0.6,y,0.07,h,barc)
    if isinstance(runs,str): runs=[(runs,WH,True)]
    text(s,0.86,y,8.35,h,runs,size=13.5,color=WH,bold=True,anchor=MSO_ANCHOR.MIDDLE,ls=1.04)

def footer(s,n,src=None):
    text(s,0.6,5.31,7.2,0.24,src or "",size=8,color=DIM)
    text(s,8.7,5.31,0.7,0.24,f"{n:02d}",size=9,color=DIM,align=PP_ALIGN.RIGHT)

def boss_avatar(s,x,y,d=0.92,label="論破王"):
    oval(s,x,y,d,d,REDBG,line=RED,line_w=2.0)
    text(s,x,y+0.10,d,0.34,"👤",size=20,color=WH,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    text(s,x,y+d-0.32,d,0.28,label,size=10,color=RED,bold=True,align=PP_ALIGN.CENTER)

def boss_bubble(s,x,y,w,h,line_text):
    """Boss attack speech bubble (red)."""
    card(s,x,y,w,h,fill=REDBG,line=RED,radius=0.12,shadow=True,line_w=1.25)
    text(s,x+0.22,y,w-0.4,h,[("ボス：", RED, True, 12), (line_text, WH, True, 13.5)],
         anchor=MSO_ANCHOR.MIDDLE,ls=1.05)

def hp_gauge(s,x,y,defeated,total=4):
    """HP pips: defeated = empty/grey, remaining = red."""
    text(s,x,y-0.02,1.5,0.26,[("BOSS HP ",DIM,True,9.5),(f"{total-defeated}/{total}",RED if defeated<total else CY,True,11)])
    px=x+1.55; pw=0.34; gap=0.10
    for i in range(total):
        c = "243154" if i<defeated else RED
        b=rect(s,px,y,pw,0.18,c);
        px+=pw+gap

# =====================================================================
# S0 — BOSS APPEARS
# =====================================================================
s=slide(BG_COVER)
text(s,0.6,0.40,8.8,0.28,"BOSS BATTLE ・ 本日の関門",size=12,color=RED,bold=True)
text(s,0.6,0.74,9.0,0.70,"ボスは「論破王」FM",size=40,color=WH,bold=True)
keymsg(s,[("論破を全部しのげたら、承認。",WH,True),("感想は要らない、ぜんぶ数字で殴り返す。",GOLD,True)],y=1.66,h=0.60,barc=RED)
# boss
boss_avatar(s,0.9,2.55,d=1.05,label="論破王 FM")
card(s,2.15,2.55,4.0,1.85,fill=REDBG,line=RED,radius=0.10,line_w=1.25)
text(s,2.4,2.70,3.55,0.30,"必殺技",size=11,color=RED,bold=True)
text(s,2.4,3.06,3.55,1.25,
     [[("「それってあなたの感想ですよね？」",WH,True,12.5)],
      [("「なんかそういうデータあるんですか？」",WH,True,12.5)],
      [("「……で、論破」",RED,True,13)]],ls=1.25)
# challenger
card(s,6.35,2.55,3.05,1.85)
text(s,6.58,2.70,2.6,0.30,"挑戦者",size=11,color=CY,bold=True)
text(s,6.58,3.06,2.65,1.25,
     [[("新人アナリスト（私）",WH,True,13)],
      [("議題＝3号ファンドから",MUT,False,11.5)],
      [("Solafuneにリード投資",CY,True,12.5)],
      [("全[仮]に数字と出典を用意",MUT,False,11)]],ls=1.2)
# rule + hp
text(s,0.6,4.62,5.7,0.34,[("勝利条件＝論破4連打（誇大／感想／倫理／なぜ今あなた達）を",MUT,False,11.5),("全て数字で返す",WH,True,11.5)])
hp_gauge(s,6.35,4.66,0,4)
footer(s,0)

# =====================================================================
# S1 — WHY US (track record)
# =====================================================================
s=slide()
header(s,"WHY US ・ なぜ私たち","なぜ私たちなのか ― 宇宙で2度、結果を出した")
keymsg(s,[("1号 2.0X → 2号 1.5X。",WH,True),("倍率が下がった“痛み”こそ、3号「集中」の根拠。",GOLD,True)])
funds=[("1号","¥30億","宇宙の裾野・部品/データ","DPI 2.0X",CY),
       ("2号","¥80億","宇宙データ・解析の成長期","DPI 1.5X",PUR),
       ("3号","¥150億","経済安保×宇宙×AIに集中","ネット3.0X 狙い",GOLD)]
x=0.6; cw=2.78
for i,(t,sz,desc,dpi,c) in enumerate(funds):
    card(s,x,2.18,cw,1.72)
    text(s,x+0.2,2.32,cw-0.4,0.34,[(t+"  ",WH,True,16),(sz,c,True,15)])
    text(s,x+0.2,2.76,cw-0.4,0.60,desc,size=11,color=MUT,ls=1.1)
    text(s,x+0.2,3.40,cw-0.4,0.40,dpi,size=17,color=c,bold=True)
    if i<2: text(s,x+cw-0.02,2.18,0.18,1.72,"▶",size=13,color=DIM,anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.CENTER)
    x+=cw+0.155
# sourcing + boss pre-empt
text(s,0.6,4.06,8.8,0.30,[("宇宙を2度回したからこそのソーシング：",MUT,False,11.5),("120カ国コミュニティ／防衛装備庁・国際機関リレーション",WH,True,11.5)])
card(s,0.6,4.46,8.8,0.62,fill=REDBG,line=RED,radius=0.10,shadow=False,line_w=1.0)
text(s,0.86,4.46,8.3,0.62,
     [("ボス：「2号で下がってるじゃないですか」　",RED,True,12),
      ("→ はい。だから分散をやめ、集中に振り切った。痛みを戦略に変えました。",WH,True,12)],
     anchor=MSO_ANCHOR.MIDDLE)
footer(s,1)

# =====================================================================
# S2 — WHY NOW (institutional changes)
# =====================================================================
s=slide()
header(s,"WHY NOW ・ なぜ今","なぜ今なのか ― 2026年、制度が動いた")
keymsg(s,[("追い風は気分じゃない。",WH,True),("輸出緩和・ファストパス・前払い",GOLD,True),("で、市場と参入条件が構造的に変わった。",WH,True)])
ch=[("① 輸出ルール緩和","協定国へ防衛装備の移転可。同志国市場が開く＝数が出て単価が下がる（量産効果）",CY),
    ("② ファストパス調達","公募〜契約が約3.5か月（従来は1年超）",GOLD),
    ("③ 前払い・部分払い","ディープテックの「死の谷」を浅くする",PUR)]
x=0.6; cw=2.78
for t,d,c in ch:
    card(s,x,2.18,cw,1.72)
    rect(s,x,2.18,cw,0.08,c)
    text(s,x+0.2,2.40,cw-0.4,0.40,t,size=13.5,color=WH,bold=True)
    text(s,x+0.2,2.90,cw-0.42,1.0,d,size=11,color=MUT,ls=1.18)
    x+=cw+0.155
text(s,0.6,4.06,8.8,0.30,[("帰結：",MUT,False,11.5),("民→官のデュアルユース（防災を入口に有事へ）が成立しやすい構造に。",WH,True,11.5)])
card(s,0.6,4.46,8.8,0.62,fill=REDBG,line=RED,radius=0.10,shadow=False,line_w=1.0)
text(s,0.86,4.46,8.3,0.62,
     [("ボス：「日本は輸出できないでしょ」　",RED,True,12),
      ("→ 2026年4月に変わりました。戦闘中の国は除外、歯止めとセットです。",WH,True,12)],
     anchor=MSO_ANCHOR.MIDDLE)
footer(s,2)

# =====================================================================
# 論破 rounds — reusable
# =====================================================================
def rebuttal(n, kicker, title, attack, counter_runs, body_paras, src=None, title_sz=21):
    s=slide()
    header(s,kicker,title,ksize=title_sz,kc=RED)
    boss_bubble(s,0.6,1.30,8.8,0.66,attack)
    keymsg(s,counter_runs,y=2.08,h=0.60,barc=CY)
    card(s,0.6,2.84,8.8,1.78)
    yy=3.00
    for para in body_paras:
        text(s,0.86,yy,8.2,0.44,para,size=11.5,color=MUT,ls=1.06)
        yy+=0.42
    # "しのいだ" + HP
    text(s,0.6,4.70,5.4,0.30,[("論破返し成功 ▶ ",CY,True,12),("ボスHP −1",WH,True,12)])
    hp_gauge(s,6.35,4.74,n,4)
    footer(s,2+n,src)
    return s

# 論破① 誇大
rebuttal(1,"論破① ・ DEBUNK ROUND 1",
    "「1.37Xで“ファンドリターナー”は盛りすぎでは？」",
    "1.37Xって全然3.0Xに足りてないですよね？",
    [("別レイヤーの話です。",WH,True),("1.37Xは“対ファンド比”、この1社のMOICは12.5X。",GOLD,True),("盛ってない、定義です。",WH,True)],
    [[("● ファンドリターナー＝1社で元手150億をほぼ1本ぶん回収＝ポジション206億／ファンド150億＝1.37X",MUT,False,11.5)],
     [("● 同じ1社を“投じた額”で見れば：206億 ÷ 総コミット16.5億 ＝ ",MUT,False,11.5),("MOIC 12.5X",GOLD,True,11.5)],
     [("● 全体3.0Xは別物：10社のべき乗則（上位2社で回収の約65%）で積む",MUT,False,11.5)],
     [("● DPI＝財布全体が何倍／MOIC＝この1社が何倍。混ぜると誤読します",DIM,False,11)]])

# 論破② 感想（EV/Sales）
rebuttal(2,"論破② ・ DEBUNK ROUND 2",
    "「EV/Sales 17倍って、あなたの感想ですよね？」",
    "その17倍、なんかそういうデータあるんですか？",
    [("感想じゃなく統計です。",WH,True),("comps中央値10倍＋要因分解(+7)、",GOLD,True),("Palantirバブルは除外済み。",WH,True)],
    [[("● 採用群 {Planet6・Synspective7・BlackSky10・Anduril14・Palantirプレ20} → 中央値10倍",MUT,False,11.5)],
     [("● 上乗せ+7倍＝ソフト比率+3／粗利優位+2／主権的信認+2 ＝ ",MUT,False,11.5),("17倍",GOLD,True,11.5)],
     [("● Palantir現況60.7倍はバブル局面として明示除外（だから保守でも15〜18倍）",MUT,False,11.5)],
     [("● 中央値10倍まで縮んでも、1社で元本の約8割を回収＝壊れない",DIM,False,11)]])

# 論破③ 倫理（Humanity Brain）
rebuttal(3,"論破③ ・ DEBUNK ROUND 3",
    "「Humanity Brainって“世論操作”の会社では？」",
    "認知戦AIって、世論を操作する危ない会社じゃないですか？",
    [("逆です。",WH,True),("用途は“認知防衛（誤情報耐性）”に限定、",GOLD,True),("攻撃的な利用は投資契約で禁止。",WH,True)],
    [[("● Humanity Brain＝認知戦対策AI・世論分布予測（防衛装備庁「社会シミュレーション」採択）",MUT,False,11.5)],
     [("● 線引き：対象国の認知“操作”ではなく、自国の認知“防衛”＝誤情報への耐性に用途限定",MUT,False,11.5)],
     [("● 投資契約に influence-ops の攻撃的利用禁止コベナンツ＋第三国移転制限",MUT,False,11.5)],
     [("● Solafune（物理GEOINT）の隣＝認知ドメインの解析層。守る側です",DIM,False,11)]],
    src="※前回ボスが唯一刺したESGの穴を、ここで塞ぐ")

# 論破④ なぜ今あなた達
rebuttal(4,"論破④ ・ DEBUNK ROUND 4",
    "「で、なぜ“今”“あなた達”が勝てるの？」",
    "別に誰でもよくないですか？なぜあなた達？",
    [("なぜ私たち × なぜ今 × なぜSolafune ",WH,True),("― 3つ揃ってます。",GOLD,True)],
    [[("● なぜ私たち：宇宙で1号2.0X・2号1.5Xの実績＝ソーシングとパターン認識を持つ",MUT,False,11.5)],
     [("● なぜ今：2026年の輸出緩和・ファストパス・前払いで参入条件が構造的に好転",MUT,False,11.5)],
     [("● なぜSolafune：ハードは大企業・ソフト/解析はスタートアップ（Anduril「生産に価値はない」）",MUT,False,11.5)],
     [("● その解析層に、二桁持分でリードできるのが我々",DIM,False,11)]])

# =====================================================================
# S_END — 決着
# =====================================================================
s=slide()
text(s,0.6,0.50,9.0,0.30,"FINISH ・ 決着",size=12,color=CY,bold=True)
text(s,0.6,1.05,9.0,0.70,"論破ポイント、全部塞ぎました",size=34,color=WH,bold=True)
keymsg(s,[("感想ゼロ・全部数字。",WH,True),("初回5.5億・持分約15%でSolafuneにリード。",GOLD,True),("Exit 206億・1.37X・ネットDPI 3.0X。",CY,True)],y=1.98,h=0.60)
# HP cleared
text(s,0.6,2.78,3.0,0.30,[("BOSS HP ",DIM,True,10),("0/4 ",CY,True,13),("撃破",CY,True,12)])
px=2.7
for i in range(4):
    rect(s,px,2.82,0.34,0.18,"243154"); px+=0.44
text(s,5.0,2.76,4.4,0.30,"誇大→定義／感想→統計／倫理→契約／なぜ今→3点回収",size=10.5,color=MUT)
# decision
card(s,0.6,3.30,8.8,0.96,fill=KMFILL,line=GOLD,radius=0.10,shadow=False)
text(s,0.86,3.42,8.3,0.40,[("投資判断：",GOLD,True,13),("初回5.5億・持分約15%でリード、リザーブ厚めでフォローオン＝ファンドリターナー候補",WH,True,12.5)])
text(s,0.86,3.84,8.3,0.36,[("次アクション：",MUT,False,11.5),("シード〜プレA単独リード組成の条件交渉、KPI連動トランシェ設計",MUT,False,11.5)])
text(s,0.6,4.45,9.0,0.40,[("「……で、論破できました？」　",RED,True,13),("地球を読む者が、次の安全保障を制する。",CY,True,13)])
footer(s,7)

prs.save(OUT)
print("saved:", OUT, "slides:", len(prs.slides._sldIdLst))
