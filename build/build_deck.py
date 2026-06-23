# -*- coding: utf-8 -*-
"""Build Sentinel Capital × Solafune pitch deck — ボス戦版（空パケ06／単線FR）.
03_スライド骨子.md の S0〜S16（17枚）に同期。
  S0 ボス登場 / S1 表紙テーゼ / S2 なぜ私たち / S3 勝ち筋(FR定義+付加価値) /
  S4 ファンド設計 / S5 なぜ今 / S6 Solafune / S7 市場SOM / S8 競争 / S9 財務 /
  S10 バリュ&Exit / S11 論破②(HP4→3) / S12 リターン&FR検算 / S13 論破①(HP3→2) /
  S14 論破③(HP2→1) / S15 リスク / S16 決着&論破④(HP1→0)
論破①②③④はボス吹き出し→数字返し＋BOSS HPゲージ（4→3→2→1→0）。
芯＝単線FR：Solafune＝唯一のファンドリターナー（突出1本）。Humanity Brain＝勝ち馬Bの実在候補（例示）。
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
BG_COVER = os.path.join(HERE, "assets", "image1.png")
BG_BODY  = os.path.join(HERE, "assets", "image13.png")
OUT = os.path.join(os.path.dirname(HERE), "outputs", "05_発表スライド.pptx")

# ---- palette（ライト＝本文白テーマ / ダーク＝扉絵・ボス演出）----
_DEF = object()
DARK = dict(CY="5BC8FF", WH="EAF1FF", MUT="9DB0D0", DIM="6E7FA6", GOLD="F5B454",
            PUR="9C8CFF", CARD="111C3C", CARDLN="2A3B66", KMFILL="13234D",
            RED="FF5C6C", REDBG="2A1530", GAUGE_OFF="243154", RULE="2C4A86")
LIGHT = dict(CY="1488C2", WH="16263F", MUT="56627A", DIM="7C8AA3", GOLD="C77D11",
             PUR="6E5AD9", CARD="EEF3FB", CARDLN="D5DEEC", KMFILL="EEF0F4",
             RED="DC3B52", REDBG="FCE9ED", GAUGE_OFF="C7D0E0", RULE="9FAEC6")

def set_theme(name):
    global CY, WH, MUT, DIM, GOLD, PUR, CARD, CARDLN, KMFILL, RED, REDBG, GAUGE_OFF, RULE
    p = LIGHT if name == "light" else DARK
    CY, WH, MUT, DIM, GOLD, PUR, CARD, CARDLN, KMFILL, RED, REDBG, GAUGE_OFF, RULE = (
        p["CY"], p["WH"], p["MUT"], p["DIM"], p["GOLD"], p["PUR"], p["CARD"],
        p["CARDLN"], p["KMFILL"], p["RED"], p["REDBG"], p["GAUGE_OFF"], p["RULE"])

set_theme("light")
FONT = "Arial"; EAFONT = "Meiryo"

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)
BLANK = prs.slide_layouts[6]

def rgb(h): return RGBColor.from_string(h)

def slide(bg=None):
    s = prs.slides.add_slide(BLANK)
    if bg:  # 扉絵・ボス演出は宇宙背景を維持
        pic = s.shapes.add_picture(bg, 0, 0, prs.slide_width, prs.slide_height)
        sp = pic._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2, sp)
    # bg=None の本文スライドは白（BLANK レイアウト既定）
    return s

def _set_font(run, size, color, bold, font=FONT):
    f = run.font
    f.size = Pt(size); f.bold = bold; f.name = font; f.color.rgb = rgb(color)
    rPr = run._r.get_or_add_rPr(); ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
    ea.set('typeface', EAFONT)

def text(s, x, y, w, h, runs, size=12, color=_DEF, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, line_spacing=1.0, font=FONT, wrap=True):
    if color is _DEF: color = MUT
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        paras = [[(runs, color, bold, size)]]
    elif runs and isinstance(runs[0], tuple):
        paras = [runs]
    else:
        paras = runs
    for pi, para in enumerate(paras):
        p = tf.paragraphs[0] if pi==0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        for tup in para:
            txt = tup[0]; c = tup[1] if len(tup)>1 and tup[1] else color
            b = tup[2] if len(tup)>2 and tup[2] is not None else bold
            sz = tup[3] if len(tup)>3 and tup[3] else size
            r = p.add_run(); r.text = txt; _set_font(r, sz, c, b, font)
    return tb

def _shadow(shape):
    spPr = shape._element.spPr
    el = spPr.makeelement(qn('a:effectLst'), {})
    sh = spPr.makeelement(qn('a:outerShdw'),
        {'blurRad':'90000','dist':'30000','dir':'5400000','rotWithShape':'0','algn':'tl'})
    clr = spPr.makeelement(qn('a:srgbClr'), {'val':'000000'})
    al = spPr.makeelement(qn('a:alpha'), {'val':'38000'})
    clr.append(al); sh.append(clr); el.append(sh); spPr.append(el)

def card(s, x, y, w, h, fill=_DEF, line=_DEF, radius=0.07, shadow=True, line_w=1.0):
    if fill is _DEF: fill = CARD
    if line is _DEF: line = CARDLN
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = rgb(fill)
    if line:
        sh.line.color.rgb = rgb(line); sh.line.width = Pt(line_w)
    else:
        sh.line.fill.background()
    try: sh.adjustments[0] = radius
    except Exception: pass
    sh.shadow.inherit = False
    if shadow: _shadow(sh)
    if sh.has_text_frame: sh.text_frame.word_wrap = True
    return sh

def bar(s, x, y, w, h, color, radius=0.18):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = rgb(color); sh.line.fill.background()
    try: sh.adjustments[0]=radius
    except Exception: pass
    sh.shadow.inherit=False
    return sh

def rect(s, x, y, w, h, color):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = rgb(color); sh.line.fill.background()
    sh.shadow.inherit=False
    return sh

def oval(s, x, y, w, h, fill, line=None, line_w=2.0):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = rgb(fill)
    if line: sh.line.color.rgb = rgb(line); sh.line.width = Pt(line_w)
    else: sh.line.fill.background()
    sh.shadow.inherit = False; return sh

def header(s, kicker, title, ksize=27, kc=_DEF):
    # アイブロウ（kicker）は廃止。タイトルを上詰めし、直下に区切り線。
    text(s, 0.6, 0.40, 8.8, 0.74, title, size=ksize, color=WH, bold=True,
         anchor=MSO_ANCHOR.MIDDLE)
    rect(s, 0.6, 1.215, 8.8, 0.018, RULE)

def keymsg(s, runs, y=1.34, h=0.62, barc=_DEF):
    if barc is _DEF: barc = CY
    card(s, 0.6, y, 8.8, h, fill=KMFILL, line=CARDLN, radius=0.06, shadow=False, line_w=0.75)
    rect(s, 0.6, y, 0.07, h, barc)
    if isinstance(runs, str): runs=[(runs, WH, True)]
    text(s, 0.86, y, 8.35, h, runs, size=13.5, color=WH, bold=True,
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.04)

def footer(s, n, src=None):
    text(s, 0.6, 5.31, 7.2, 0.24, src or "", size=8, color=DIM)
    text(s, 8.7, 5.31, 0.7, 0.24, f"{n:02d}", size=9, color=DIM, align=PP_ALIGN.RIGHT)

def metric(s, x, y, w, h, big, big_c, label):
    card(s, x, y, w, h)
    text(s, x+0.18, y+0.16, w-0.36, 0.62, big, size=30, color=big_c, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x+0.18, y+h-0.42, w-0.36, 0.34, label, size=11, color=MUT, anchor=MSO_ANCHOR.MIDDLE)

# ---- boss helpers ----
def boss_avatar(s, x, y, d=0.92, label="論破王 FM"):
    oval(s, x, y, d, d, REDBG, line=RED, line_w=2.0)
    text(s, x, y+0.10, d, 0.34, "👤", size=20, color=WH, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x, y+d-0.32, d, 0.28, label, size=10, color=RED, bold=True, align=PP_ALIGN.CENTER)

def boss_bubble(s, x, y, w, h, line_text):
    card(s, x, y, w, h, fill=REDBG, line=RED, radius=0.12, shadow=True, line_w=1.25)
    text(s, x+0.22, y, w-0.4, h, [("ボス：", RED, True, 12), (line_text, WH, True, 13.5)],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)

def hp_gauge(s, x, y, defeated, total=4):
    text(s, x, y-0.02, 1.5, 0.26, [("BOSS HP ", DIM, True, 9.5),
         (f"{total-defeated}/{total}", RED if defeated < total else CY, True, 11)])
    px = x+1.55; pw = 0.34; gap = 0.10
    for i in range(total):
        c = GAUGE_OFF if i < defeated else RED
        rect(s, px, y, pw, 0.18, c); px += pw+gap

def rebuttal(page, defeated, kicker, title, attack, counter_runs, body_paras, src=None, title_sz=20):
    s = slide()
    header(s, kicker, title, ksize=title_sz, kc=RED)
    boss_bubble(s, 0.6, 1.30, 8.8, 0.66, attack)
    keymsg(s, counter_runs, y=2.08, h=0.60, barc=CY)
    card(s, 0.6, 2.84, 8.8, 1.78)
    yy = 3.00
    for para in body_paras:
        text(s, 0.86, yy, 8.2, 0.44, para, size=11.5, color=MUT, line_spacing=1.06)
        yy += 0.42
    text(s, 0.6, 4.70, 5.4, 0.30, [("論破返し成功 ▶ ", CY, True, 12), ("ボスHP −1", WH, True, 12)])
    hp_gauge(s, 6.35, 4.74, defeated, 4)
    footer(s, page, src)
    return s

def appendix(ap_no, title, lines, src=None, fsz=9.5, lh=0.265, cpl=48):
    """補足（QA用・本編では飛ばす・0秒）スライド。lines＝('H',見出し)／文字列／runsリストの混在。
    行幅を概算して折返し行数ぶん自動で送る（密な台帳でも重なり/はみ出しを防ぐ）。"""
    s = slide()
    text(s, 0.6, 0.34, 6.0, 0.26, "APPENDIX ・ 補足（本編では飛ばす｜QA用）", size=11, color=GOLD, bold=True)
    text(s, 0.6, 0.62, 7.1, 0.52, f"AP{ap_no}　{title}", size=18, color=WH, bold=True,
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, 7.0, 0.40, 2.4, 0.30, "0秒", size=10, color=DIM, bold=True,
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, 0.6, 1.205, 8.8, 0.016, RULE)
    card(s, 0.6, 1.34, 8.8, 3.82, shadow=False)
    yy = 1.50
    for ln in lines:
        if isinstance(ln, tuple) and len(ln) == 2 and ln[0] == "H":
            text(s, 0.84, yy+0.03, 8.4, 0.26, ln[1], size=11, color=CY, bold=True)
            yy += 0.31
        else:
            body = [(ln, MUT, False, fsz)] if isinstance(ln, str) else list(ln)
            raw = ln if isinstance(ln, str) else "".join(t[0] for t in ln)
            runs = [("▸ ", GOLD, False, fsz)] + body
            w = sum(1.0 if ord(c) > 0x2E80 else 0.55 for c in raw) + 2
            nl = max(1, int((w + cpl - 1) // cpl))
            text(s, 0.92, yy, 8.42, lh*nl, runs, size=fsz, color=MUT, line_spacing=1.04)
            yy += lh*nl + 0.03
    text(s, 0.6, 5.33, 7.6, 0.24, src or "", size=8, color=DIM)
    text(s, 8.5, 5.33, 0.9, 0.24, f"AP{ap_no} / 9", size=9, color=DIM, align=PP_ALIGN.RIGHT)
    return s

# =====================================================================
# S0 — BOSS APPEARS（つかみ・論破王FM登場 / HP 4/4）
# =====================================================================
set_theme("dark")
s = slide(BG_COVER)
text(s, 0.6, 0.40, 8.8, 0.28, "BOSS BATTLE ・ 本日の投資委員会", size=12, color=RED, bold=True)
text(s, 0.6, 0.74, 9.0, 0.70, "ボスは「論破王」FM", size=40, color=WH, bold=True)
keymsg(s, [("論破を全部しのげたら、承認。", WH, True),
           ("感想は要らない、ぜんぶ数字で殴り返す。", GOLD, True)], y=1.66, h=0.60, barc=RED)
boss_avatar(s, 0.9, 2.55, d=1.05, label="論破王 FM")
card(s, 2.15, 2.55, 4.0, 1.85, fill=REDBG, line=RED, radius=0.10, line_w=1.25)
text(s, 2.4, 2.70, 3.55, 0.30, "必殺技", size=11, color=RED, bold=True)
text(s, 2.4, 3.06, 3.55, 1.25,
     [[("「それってあなたの感想ですよね？」", WH, True, 12.5)],
      [("「なんかそういうデータあるんですか？」", WH, True, 12.5)],
      [("「……で、論破」", RED, True, 13)]], line_spacing=1.25)
card(s, 6.35, 2.55, 3.05, 1.85)
text(s, 6.58, 2.70, 2.6, 0.30, "挑戦者", size=11, color=CY, bold=True)
text(s, 6.58, 3.06, 2.65, 1.25,
     [[("新人アナリスト（私）", WH, True, 13)],
      [("議題＝3号ファンドから", MUT, False, 11.5)],
      [("Solafuneにリード投資", CY, True, 12.5)],
      [("全[仮]に数字と出典を用意", MUT, False, 11)]], line_spacing=1.2)
text(s, 0.6, 4.62, 5.7, 0.34, [("勝利条件＝論破4連打（誇大／感想／倫理／なぜ今あなた達）を", MUT, False, 11.5),
     ("全て数字で返す", WH, True, 11.5)])
hp_gauge(s, 6.35, 4.66, 0, 4)
footer(s, 0)

# =====================================================================
# A — アジェンダ（S0の次・プレゼンでは飛ばす＝0秒）／以降は白テーマ
# =====================================================================
set_theme("light")
s = slide()
header(s, "", "アジェンダ")
# 端に小さく「プレゼンでは飛ばす（0秒）／詳細はAppendix」＝PDFで読む採点者向けの道筋ページ
text(s, 5.4, 0.46, 4.0, 0.30, "プレゼンでは飛ばす（0秒）／詳細はAppendix", size=9.5, color=DIM, bold=True,
     align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
keymsg(s, [("本日の論証は一本道＝「集中で勝つファンドが、なぜ今Solafuneに賭けるか」を、", WH, True),
           ("①設計→②対象→③検算・論破→④決着", GOLD, True),
           ("の順に全部数字で示す。", WH, True)])
ag_items = [("①", "ファンド設計｜なぜ我々が・どう勝つか", "トラックレコード→勝ち筋→ファンド設計", "S1–S4"),
            ("②", "投資対象｜なぜ今・なぜSolafune", "制度変化→OS→市場→競合→財務→Exit", "S5–S10"),
            ("③", "検算・論破｜誇大/感想/倫理を撃破", "リターン検算→論破①②③→リスク", "S11–S15"),
            ("④", "決着｜なぜ今・あなた達＝総回収", "投資判断（論破④で締める）", "S16")]
yy = 2.18; rh = 0.62; gap = 0.10
for num, label, sub, rng in ag_items:
    card(s, 0.6, yy, 8.8, rh)
    d = 0.42
    oval(s, 0.84, yy+(rh-d)/2, d, d, CY)
    text(s, 0.84, yy+(rh-d)/2, d, d, num, size=14, color="FFFFFF", bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 1.54, yy+0.08, 6.1, 0.30, label, size=13.5, color=WH, bold=True)
    text(s, 1.54, yy+0.37, 6.1, 0.24, sub, size=9.5, color=MUT)
    text(s, 7.9, yy, 1.3, rh, rng, size=12.5, color=GOLD, bold=True,
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    yy += rh + gap
text(s, 0.6, 5.04, 8.8, 0.26,
     "※複雑な導出・詳細データは Appendix（AP1〜9）に退避＝本編では飛ばし、QAで開く。",
     size=9.5, color=DIM)

# =====================================================================
# A2 — エグゼクティブサマリー
# =====================================================================
s = slide()
header(s, "", "エグゼクティブサマリー ＝ 逆算マップ")
# 結論カード（先出し）
card(s, 0.6, 1.34, 8.8, 0.56, fill=KMFILL, line=CARDLN, radius=0.06, shadow=False, line_w=0.75)
rect(s, 0.6, 1.34, 0.07, 0.56, CY)
text(s, 0.86, 1.34, 8.35, 0.56,
     [("結論：3号150億から初回5.5億・持分約15%で、", WH, True), ("Solafune にリード投資する。", GOLD, True)],
     size=13.5, color=WH, bold=True, anchor=MSO_ANCHOR.MIDDLE)
# 逆算の予告ラベル
text(s, 0.6, 1.98, 8.8, 0.26,
     [("本日の証明は一本の逆算 ― 目標 DPI3.0X から各リンクが一意に決まる", WH, True, 11),
      ("（数字下＝証明スライド）", MUT, False, 10)])
# 逆算カスケード（6ノードを ← で連結＝目標から下流へ逆算／S12で閉じる）
es_nodes = [("3.0X", "ネットDPI", "S12", PUR), ("¥206億", "1社FR", "S3·S12", GOLD),
            ("¥1,717億", "Exit", "S10", GOLD), ("17倍", "EV/Sales", "S10", CY),
            ("¥90億≒SOM", "必要売上=市場", "S7", CY), ("15→12%", "seed持分", "S4·S6", GOLD)]
nx = 0.6; nbw = 1.28; ngap = 0.20; ny = 2.34; nch = 1.28
for i, (big, lab, ref, c) in enumerate(es_nodes):
    card(s, nx, ny, nbw, nch)
    text(s, nx+0.04, ny+0.12, nbw-0.08, 0.44, big, size=13, color=c, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, nx+0.04, ny+0.58, nbw-0.08, 0.36, lab, size=8.7, color=MUT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    text(s, nx+0.04, ny+0.98, nbw-0.08, 0.24, ref, size=8.5, color=CY, bold=True, align=PP_ALIGN.CENTER)
    if i < len(es_nodes)-1:
        text(s, nx+nbw, ny, ngap, nch, "←", size=15, color=DIM, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    nx += nbw + ngap
# 3つの「なぜ」（鎖を成立させる前提）
es_rows = [("なぜ今", "2026年の輸出緩和・調達ファストパス・前払いで参入条件が構造変化"),
           ("なぜ我々", [[("・宇宙VCを2回運用した実績（目利き）", WH, False, 8.8)],
                        [("・取締役・120カ国網・政府パイプで一緒に伸ばす", WH, False, 8.8)]]),
           ("なぜSolafune", "衛星を持たない解析特化＝二桁持分を保てる希少企業（FR候補）")]
ex = 0.6; ecw = 2.86; egap = 0.11
for lbl, desc in es_rows:
    card(s, ex, 3.78, ecw, 0.76)
    text(s, ex+0.16, 3.84, ecw-0.32, 0.26, lbl, size=11.5, color=CY, bold=True)
    text(s, ex+0.16, 4.12, ecw-0.32, 0.42, desc, size=8.8, color=WH, line_spacing=1.04)
    ex += ecw + egap
text(s, 0.6, 4.70, 8.8, 0.30,
     [("この鎖を8分で1リンクずつ実数で証明し、", MUT, False, 10.5),
      ("S12で閉じる（答え合わせ）", GOLD, True, 10.5),
      ("。残リスクは契約条項で先回り。", MUT, False, 10.5)])

# =====================================================================
# S1 — 表紙＋投資テーゼ（ES の後・本編の扉）
# =====================================================================
set_theme("dark")
s = slide(BG_COVER)
text(s, 0.6, 0.50, 8.0, 0.30, "VC FUND PITCH ・ 投資テーゼ", size=12, color=CY, bold=True)
text(s, 0.6, 1.05, 9.0, 0.85, "SENTINEL CAPITAL", size=44, color=WH, bold=True)
text(s, 0.6, 1.78, 9.0, 0.70, "PARTNERS  ×  Solafune", size=30, color=CY, bold=True)
keymsg(s, [("3号150億から初回5.5億・エントリー持分約15%で、", WH, True),
           ("Solafuneにリード投資する。", GOLD, True)], y=2.58, h=0.58)
text(s, 0.6, 3.40, 8.9, 0.62,
     [("テーゼ：衛星を持たない", MUT, False, 13),
      ("「Planetary Intelligence OS」＝解析・意思決定レイヤー特化", WH, True, 13),
      ("の希少な日本のデュアルユース・ソフト企業。", MUT, False, 13)], line_spacing=1.15)
# mini summary badges
badges = [("3号 ¥150億", CY), ("集中 10–12社", GOLD), ("ネットDPI 3.0X", PUR)]
x = 0.6
for b, c in badges:
    card(s, x, 4.20, 2.55, 0.58, radius=0.16)
    text(s, x, 4.20, 2.55, 0.58, b, size=15, color=c, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += 2.73
# 運用主体（誰が運用するか＝設定を最初に明記／評価基準①）
text(s, 0.6, 4.92, 8.9, 0.32,
     [("運用：Sentinel Capital Partners｜経済安保特化GP", WH, True, 10),
      ("（宇宙2号運用＋安保クリアランス・防衛/商社オペレーティングパートナー）", MUT, False, 9.5)])
footer(s, 1, "→詳細：S3・AP9（FR4条件・チーム/付加価値）")

# =====================================================================
# S2 — なぜ私たち（トラックレコード）／以降は白テーマ
# =====================================================================
set_theme("light")
s = slide()
header(s, "WHY US ・ なぜ私たち", "なぜ私たちなのか ― 宇宙で2度、結果を出した")
keymsg(s, [("1号 2.0X → 2号 1.5X。", WH, True),
           ("倍率が下がった“痛み”こそ、3号「集中」の根拠。", GOLD, True)])
funds = [("1号", "¥30億", "宇宙の裾野・部品/データ", "DPI 2.0X", CY),
         ("2号", "¥80億", "宇宙データ・解析の成長期", "DPI 1.5X", PUR),
         ("3号", "¥150億", "経済安保×宇宙×AIに集中", "ネット3.0X 狙い", GOLD)]
x = 0.6; cw = 2.78
for i, (t, sz, desc, dpi, c) in enumerate(funds):
    card(s, x, 2.18, cw, 1.72)
    text(s, x+0.2, 2.32, cw-0.4, 0.34, [(t+"  ", WH, True, 16), (sz, c, True, 15)])
    text(s, x+0.2, 2.76, cw-0.4, 0.60, desc, size=11, color=MUT, line_spacing=1.1)
    text(s, x+0.2, 3.40, cw-0.4, 0.40, dpi, size=17, color=c, bold=True)
    if i < 2: text(s, x+cw-0.02, 2.18, 0.18, 1.72, "▶", size=13, color=DIM, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    x += cw+0.155
text(s, 0.6, 4.06, 8.8, 0.30, [("宇宙を2度回したからこそのソーシング：", MUT, False, 11.5),
     ("120カ国コミュニティ／防衛装備庁・国際機関リレーション", WH, True, 11.5)])
card(s, 0.6, 4.46, 8.8, 0.62, fill=REDBG, line=RED, radius=0.10, shadow=False, line_w=1.0)
text(s, 0.86, 4.46, 8.3, 0.62,
     [("ボス：「2号で下がってるじゃないですか」　", RED, True, 12),
      ("→ はい。だから分散をやめ、集中に振り切った。痛みを戦略に変えました。", WH, True, 12)],
     anchor=MSO_ANCHOR.MIDDLE)
footer(s, 2, "→詳細：AP9（1号2号の中身・ソーシング機構・“人の流れ”思想）")

# =====================================================================
# S3 — 勝ち筋（べき乗則＋FR定義＋我々の付加価値）
# =====================================================================
s = slide()
header(s, "HOW VC WINS ・ 勝ち筋", "1社で元手を返す“ファンドリターナー”に張る")
keymsg(s, [("ネット3.0Xは“全社そこそこ”では作れない——", WH, True),
           ("少数のアウトライヤーを必ず1社仕込む", GOLD, True)])
# power-law bars (left)
card(s, 0.6, 2.18, 4.55, 2.42)
text(s, 0.85, 2.32, 4.05, 0.30, "10社のリターン分布（イメージ）", size=11.5, color=CY, bold=True)
heights = [(1.45, GOLD, "FR"), (0.98, CY, "7-8X"), (0.66, MUT, ""), (0.58, MUT, ""), (0.38, DIM, ""),
           (0.30, DIM, ""), (0.22, DIM, ""), (0.16, DIM, ""), (0.12, DIM, ""), (0.08, DIM, "")]
bx = 0.95; bw = 0.30; base = 4.35
for hgt, c, lab in heights:
    bar(s, bx, base-hgt, bw, hgt, c)
    if lab: text(s, bx-0.07, base-hgt-0.22, bw+0.4, 0.20, lab, size=8.5, color=c, bold=True)
    bx += 0.40
text(s, 0.85, 4.38, 4.05, 0.22, "突出1本（Solafune＝FR）が回収の大半を作る", size=10.5, color=MUT)
# FR 4 conditions (right)
card(s, 5.35, 2.18, 4.05, 2.42)
text(s, 5.58, 2.32, 3.6, 0.30, "ファンドリターナーの4条件（Solafune＝全○）", size=11, color=CY, bold=True)
conds = [("①", "1社で150億超のポジションが射程（高い天井）"),
         ("②", "シード〜Aでリードしエグジット持分10〜12%確保"),
         ("③", "日本のIPOか国内プライム・商社M&Aで回収"),
         ("④", "センサー非依存の解析レイヤーという自社エッジ")]
yy = 2.68
for num, desc in conds:
    text(s, 5.58, yy, 0.35, 0.40, num, size=13, color=GOLD, bold=True)
    text(s, 5.95, yy, 3.30, 0.44, desc, size=10.5, color=MUT, line_spacing=1.02)
    yy += 0.44
# 付加価値 strip
card(s, 0.6, 4.74, 8.8, 0.50, fill=KMFILL, line=None, radius=0.12, shadow=False)
text(s, 0.86, 4.74, 8.3, 0.50,
     [("我々の付加価値（＝当たり馬を勝たせる）：", GOLD, True, 10.5),
      ("取締役就任／120カ国コミュニティ／政府リレーション／フォローオン　＝", WH, True, 10.5),
      ("起業家>投資家で伴走", CY, True, 10.5)], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 3, "→詳細：AP9（FR4条件・付加価値4点・deal access4点の根拠）")

# =====================================================================
# S4 — ファンド設計（件数／チケット／持分／算数）
# =====================================================================
s = slide()
header(s, "FUND DESIGN ・ ファンド設計", "数字で語るファンドの骨格")
keymsg(s, [("150億・集中10社・初回5.5億で持分", WH, True), ("約15%", GOLD, True),
           ("、狙うはネット", WH, True), ("DPI 3.0X", CY, True)])
mw, mh, gx = 2.78, 1.18, 2.93
cells = [("¥150億", CY, "ファンド総額（3号）"), ("3.0X", GOLD, "目標ネットDPI"),
         ("10–12社", CY, "集中ポートフォリオ"), ("約15%", GOLD, "エントリー持分（リード）"),
         ("2 / 20", CY, "管理報酬 / 成功報酬"), ("10年", GOLD, "運用期間")]
for i, (b, c, l) in enumerate(cells):
    col = i % 3; row = i // 3
    metric(s, 0.6+col*gx, 2.18+row*1.30, mw, mh, b, c, l)
card(s, 0.6, 4.82, 8.8, 0.52, fill=KMFILL, line=None, radius=0.12, shadow=False)
text(s, 0.86, 4.82, 8.3, 0.52,
     [("算数：初回5.5億 ÷ 投資後36.5億（Pre 31＋5.5）＝ ", MUT, False, 12),
      ("約15%", GOLD, True, 12),
      ("。持分はチケット増でなくPre交渉で確保。チーム ", MUT, False, 12),
      ("GP2＋アソシ2＋クリアランス保持", WH, True, 12),
      ("＝件数の天井10〜12社。", MUT, False, 12)], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 4, "→詳細：AP4（LP構成40/30/30の合意ロジック・件数算術・Pre逆算）")

# =====================================================================
# S5 — なぜ今（制度変化）
# =====================================================================
s = slide()
header(s, "WHY NOW ・ なぜ今", "なぜ今なのか ― 2026年、制度が動いた")
keymsg(s, [("追い風は気分じゃない。", WH, True),
           ("輸出緩和・ファストパス・前払い", GOLD, True),
           ("で、市場と参入条件が構造的に変わった。", WH, True)])
ch = [("① 輸出ルール緩和", "協定国へ防衛装備の移転可。同志国市場が開く＝数が出て単価が下がる（量産効果）", CY),
      ("② ファストパス調達", "公募〜契約が約3.5か月（従来は1年超）", GOLD),
      ("③ 前払い・部分払い", "ディープテックの「死の谷」を浅くする", PUR)]
x = 0.6; cw = 2.78
for t, d, c in ch:
    card(s, x, 2.10, cw, 1.62)
    rect(s, x, 2.10, cw, 0.08, c)
    text(s, x+0.2, 2.30, cw-0.4, 0.40, t, size=13, color=WH, bold=True)
    text(s, x+0.2, 2.78, cw-0.42, 0.92, d, size=10.5, color=MUT, line_spacing=1.14)
    x += cw+0.155
text(s, 0.6, 3.92, 8.8, 0.46,
     [("帰結＝民→官のデュアルユースが成立しやすい構造に。公的需要も実在：", MUT, False, 11),
      ("防衛費8.7兆／宇宙戦略基金10年1兆／衛星コンステPFI 2,832億。", WH, True, 11)], line_spacing=1.12)
card(s, 0.6, 4.46, 8.8, 0.62, fill=REDBG, line=RED, radius=0.10, shadow=False, line_w=1.0)
text(s, 0.86, 4.46, 8.3, 0.62,
     [("ボス：「日本は輸出できないでしょ」　", RED, True, 12),
      ("→ 2026年4月に変わりました。戦闘中の国は除外、歯止めとセットです。", WH, True, 12)],
     anchor=MSO_ANCHOR.MIDDLE)
footer(s, 5, "出典：01前提パック §6 ／ →詳細：AP6（公的需要プールの金額：防衛費8.7兆／PFI2,831億／基金1兆 ほか）")

# =====================================================================
# S6 — Solafuneとは（Planetary Intelligence OS・moat）
# =====================================================================
s = slide()
header(s, "THE TARGET ・ 投資対象", "Solafune ― 地球を読むOS")
keymsg(s, [("衛星画像・公開情報・通信という様々な“目”を束ねるAI＝", WH, True),
           ("Planetary Intelligence OS", CY, True)])
srcs = [("GEOINT", "衛星・地理空間"), ("OSINT", "公開情報・Web"), ("SIGINT", "通信・電波")]
yy = 2.22
for t, d in srcs:
    card(s, 0.6, yy, 2.2, 0.64, radius=0.14)
    text(s, 0.8, yy+0.04, 1.9, 0.30, t, size=13.5, color=WH, bold=True)
    text(s, 0.8, yy+0.36, 1.9, 0.24, d, size=9.5, color=MUT)
    yy += 0.74
card(s, 3.15, 2.34, 1.95, 1.78, fill="0E2A52", line=CY, radius=0.5)
text(s, 3.15, 2.56, 1.95, 0.40, "Intelligence", size=13, color="EAF1FF", bold=True, align=PP_ALIGN.CENTER)
text(s, 3.15, 2.96, 1.95, 0.40, "OS", size=20, color="5BC8FF", bold=True, align=PP_ALIGN.CENTER)
uses = ["重要鉱物の探査・違法採掘モニタリング（SHIGEN AI）",
        "防衛省・警察庁・内閣府 ＋ アフリカ諸国・国連機関",
        "120カ国超の開発者コミュニティが解析を供給"]
yy = 2.22
for u in uses:
    card(s, 5.45, yy, 3.95, 0.64, radius=0.14)
    text(s, 5.68, yy, 3.55, 0.64, [("▸ ", CY, True, 11.5), (u, MUT, False, 11)], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    yy += 0.74
# moat strip
card(s, 0.6, 4.50, 8.8, 0.56, fill=KMFILL, line=None, radius=0.12, shadow=False)
text(s, 0.86, 4.50, 8.3, 0.56,
     [("moat：", CY, True, 10.5), ("①センサー非依存のマルチソース解析　②120カ国コミュニティ　③政府・国際機関リレーション　④自然言語UI", WH, False, 10.5),
      ("　→ 資本効率高く二桁持分を保ちやすい", MUT, False, 10.5)], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 6, "→詳細：AP5（moat4点の反証条件・受注/MoU実績）")

# =====================================================================
# S7 — 市場（TAM/SAM/SOM・ボトムアップ）
# =====================================================================
s = slide()
header(s, "MARKET ・ 市場規模（SOM）", "“誰に・いくらで・何件”で積む")
keymsg(s, [("2032年SOM ", WH, True), ("約101億円", GOLD, True),
           (" は、Exit逆算の必要売上 約88億と一致＝市場と財務が一本の線", WH, True)])
segs = [("A 日本政府", "12件×2.5億", 30, CY), ("B 宇宙基金", "5件×3億", 15, PUR),
        ("C 同志国MoU", "8×50%×1.5億", 6, MUT), ("D 経済安保", "SHIGEN AI", 14, GOLD),
        ("E 民間", "資源・インフラ監視", 20, CY)]
text(s, 0.6, 2.16, 8.0, 0.28, "2030年 SOM 積み上げ ＝ 85億円 / 38件（TAM：GeoAI 2030年629億ドル）", size=11.5, color=CY, bold=True)
x = 0.6; scale = 8.8/85.0
for name, unit, v, c in segs:
    w = v*scale
    bar(s, x, 2.50, w-0.04, 0.46, c, radius=0.10)
    text(s, x, 2.50, w-0.04, 0.46, f"{v}", size=12, color="0B1220", bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += w
x = 0.6; cw = 1.70
for name, unit, v, c in segs:
    card(s, x, 3.18, cw-0.06, 1.05)
    rect(s, x, 3.18, cw-0.06, 0.07, c)
    text(s, x+0.12, 3.32, cw-0.26, 0.34, name, size=11, color=WH, bold=True)
    text(s, x+0.12, 3.66, cw-0.26, 0.30, unit, size=8.7, color=MUT, line_spacing=1.0)
    text(s, x+0.12, 3.96, cw-0.26, 0.24, f"¥{v}億", size=12.5, color=c, bold=True)
    x += cw
card(s, 0.6, 4.40, 8.8, 0.62, fill=KMFILL, line=None, radius=0.12, shadow=False)
text(s, 0.86, 4.40, 8.3, 0.62,
     [("2030年 ", MUT, False, 12), ("85億/38件", WH, True, 12),
      (" → 2032年 ", MUT, False, 12), ("101億/46件", GOLD, True, 12),
      ("。シェア仮定なし、単価×件数の積み上げ［C転換率が最も要実証］", MUT, False, 11)],
     anchor=MSO_ANCHOR.MIDDLE)
footer(s, 7, "出典：01前提パック §2-3 ／ →詳細：AP2（TAMレイヤー・SAM充当率・SOM5セグメント表・感度81〜121億）")

# =====================================================================
# S8 — 競争構造とポジショニング
# =====================================================================
s = slide()
header(s, "COMPETITION ・ 競争構造", "競合の真正面でなく“その上”に立つ")
keymsg(s, [("ハードは大企業・ソフト/解析はスタートアップ。", WH, True),
           ("Solafuneはマルチソース解析レイヤーで独自ポジション", GOLD, True)])
card(s, 0.6, 2.18, 4.30, 2.30)
text(s, 0.85, 2.32, 3.9, 0.30, "グローバル", size=12, color=CY, bold=True)
text(s, 0.85, 2.72, 3.85, 1.7,
     [[("● 防衛/解析テック：Palantir・Anduril・Helsing・Shield AI", MUT, False, 11)],
      [("● 防衛テックVC投資は2025年に過去最高 491億ドル", MUT, False, 11)],
      [("→ 解析・統合レイヤーが主戦場化", CY, True, 11)]], line_spacing=1.3)
card(s, 5.1, 2.18, 4.30, 2.30)
text(s, 5.35, 2.32, 3.9, 0.30, "国内（＝買い手候補でもある）", size=12, color=GOLD, bold=True)
text(s, 5.35, 2.72, 3.85, 1.7,
     [[("● Synspective＝上場で1社1,500億Exitを実証", MUT, False, 11)],
      [("● Axelspace・iQPS・重工プライム＝M&A出口の買い手", MUT, False, 11)],
      [("→ Solafuneは各社の“上”の解析層に乗る", WH, True, 11)]], line_spacing=1.3)
card(s, 0.6, 4.58, 8.8, 0.48, fill=KMFILL, line=None, radius=0.12, shadow=False)
text(s, 0.86, 4.58, 8.3, 0.48,
     [("Anduril「生産そのものに価値はない」＝", MUT, False, 11),
      ("Solafuneはその解析層に日本政府アンカーで先行（→論破④で回収）", WH, True, 11)],
     anchor=MSO_ANCHOR.MIDDLE)
footer(s, 8, "→詳細：AP1（グローバル/国内compsの評価額・倍率一覧）")

# =====================================================================
# S9 — 財務（5年売上・GM）
# =====================================================================
s = slide()
header(s, "QUANT ・ 定量モデル", "5年で描く成長（仮・要検証）")
keymsg(s, [("売上 4億→85億（CAGR約113%）×粗利 ", WH, True), ("60→72%", GOLD, True),
           ("。KPI閾値クリアで段階投資", WH, True)])
card(s, 0.6, 2.16, 5.7, 2.88)
text(s, 0.85, 2.30, 5.2, 0.28, "売上高の推移（億円）", size=11.5, color=CY, bold=True)
revs = [("26", 4), ("27", 12), ("28", 30), ("29", 55), ("30", 85)]
maxr = 85; bx = 1.05; bw = 0.72; base = 4.62; gap = 0.32
for yr, v in revs:
    h = (v/maxr)*1.9
    bar(s, bx, base-h, bw, h, CY if yr != "30" else GOLD)
    text(s, bx-0.1, base-h-0.24, bw+0.2, 0.22, str(v), size=10, color=WH, bold=True, align=PP_ALIGN.CENTER)
    text(s, bx-0.1, base+0.04, bw+0.2, 0.22, "'"+yr, size=9.5, color=MUT, align=PP_ALIGN.CENTER)
    bx += bw+gap
text(s, 0.85, 4.74, 5.2, 0.24, "日本政府 → 同志国 → 民間（グローバルサウス）の段階展開", size=10, color=MUT)
rights = [("60% → 72%", GOLD, "グロスマージン（ソフト/解析）"),
          ("¥31→36.5億", CY, "投資前Pre→投資後Post-Val"),
          ("年率 約2倍", PUR, "政府リカーリングのフォローKPI")]
yy = 2.16
for b, c, l in rights:
    card(s, 6.55, yy, 2.85, 0.90)
    text(s, 6.75, yy+0.08, 2.5, 0.44, b, size=19, color=c, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 6.77, yy+0.54, 2.5, 0.30, l, size=9.7, color=MUT)
    yy += 0.99
footer(s, 9, "出典：02財務モデル（PL/前提シート） ／ →詳細：AP2（売上セグメント内訳）")

# =====================================================================
# S10 — バリュエーション＆Exit（comps・1,717億）
# =====================================================================
s = slide()
header(s, "EXIT ・ バリュエーション", "1社でファンドを返す")
keymsg(s, [("Exit売上101億 × ", WH, True), ("EV/Sales 17倍", CY, True),
           (" → 企業価値 ", WH, True), ("約1,717億円", GOLD, True),
           ("（持分12%で206億）", WH, True)])
steps = [("売上 ¥101億", "2032E", CY), ("× 17倍", "EV/Sales（保守15–18）", MUT),
         ("¥1,717億", "Exit企業価値", GOLD), ("× 12%", "Exit持分", MUT), ("¥206億", "ポジション", PUR)]
x = 0.6; cw = 1.62; gap = 0.18
for i, (b, l, c) in enumerate(steps):
    card(s, x, 2.16, cw, 1.12)
    text(s, x+0.1, 2.30, cw-0.2, 0.50, b, size=(15 if "¥" in b and "億" in b else 16), color=c, bold=True, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    text(s, x+0.1, 2.84, cw-0.2, 0.34, l, size=9.3, color=MUT, align=PP_ALIGN.CENTER, line_spacing=1.0)
    if i < len(steps)-1:
        text(s, x+cw-0.02, 2.16, gap+0.1, 1.12, "▸", size=14, color=DIM, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    x += cw+gap
card(s, 0.6, 3.50, 5.55, 1.55)
text(s, 0.85, 3.64, 5.0, 0.30, "EV/Sales 17倍の統計的導出", size=11.5, color=CY, bold=True)
text(s, 0.85, 4.00, 5.35, 1.0,
     [[("“国しか客になれない国産インフラ”＝普通のソフトの約1.7倍の値札", CY, False, 10.5)],
      [("comps中央値10倍 ＋ ソフト+3／粗利+2／主権+2 ＝ ", MUT, False, 10.5), ("17倍", GOLD, True, 10.5)],
      [("Palantirプレ20・現況60.7倍はバブルとして除外", DIM, False, 10)]], line_spacing=1.18)
card(s, 6.35, 3.50, 3.05, 1.55, fill=KMFILL, line=GOLD)
text(s, 6.58, 3.66, 2.6, 0.30, "ファンドリターナー成立", size=11.5, color=GOLD, bold=True)
text(s, 6.58, 4.02, 2.6, 0.5, [("¥206億", WH, True, 22), (" ÷ ¥150億", MUT, False, 12)], anchor=MSO_ANCHOR.MIDDLE)
text(s, 6.58, 4.56, 2.7, 0.40, [("＝ 対ファンド ", MUT, False, 11), ("1.37X", CY, True, 13), ("／総コミMOIC ", MUT, False, 11), ("12.5X", GOLD, True, 13)], line_spacing=1.0)
footer(s, 10, "出典：01 §5・02 ／ →詳細：AP1（採用群5社の倍率・統計・要因分解+3/+2/+2・バブル除外根拠）")

# =====================================================================
# S11 — 論破②「EV/Sales 17倍＝感想ですよね？」（HP4→3）
# =====================================================================
rebuttal(11, 1, "論破② ・ DEBUNK ROUND 1（感想）",
    "「EV/Sales 17倍って、あなたの感想ですよね？」",
    "その17倍、なんかそういうデータあるんですか？",
    [("感想じゃなく統計です。", WH, True), ("comps中央値10倍＋要因分解(+7)、", GOLD, True), ("Palantirバブルは除外済み。", WH, True)],
    [[("● 採用群 {Planet6・Synspective7・BlackSky10・Anduril14・Palantirプレ20} → 中央値10倍", MUT, False, 11.5)],
     [("● 上乗せ+7倍＝ソフト比率+3／粗利優位+2／主権的信認+2 ＝ ", MUT, False, 11.5), ("17倍", GOLD, True, 11.5)],
     [("● Palantir現況60.7倍はバブル局面として明示除外（だから保守でも15〜18倍）", MUT, False, 11.5)],
     [("● 中央値10倍まで縮んでも、1社で元本の約8割を回収＝壊れない", DIM, False, 11)]],
    src="→詳細：AP1（採用群5社・要因分解+3/+2/+2・感度12〜20倍）")

# =====================================================================
# S12 — リターン＆ファンドリターナー検算＆感度
# =====================================================================
s = slide()
header(s, "RETURN ・ 逆算カスケード回収", "逆算が閉じる ― Solafune206億で回収")
keymsg(s, [("ESの逆算がここで閉じる ＝ ", WH, True),
           ("3.0X←206億←1,717億←17倍←売上101億≒SOM←seed15%", GOLD, True),
           ("（飛躍ゼロ）", WH, True)])
# big number left
card(s, 0.6, 2.18, 3.05, 1.95, fill=KMFILL, line=GOLD)
text(s, 0.8, 2.46, 2.7, 0.70, "¥206億", size=34, color=WH, bold=True, anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.8, 3.34, 2.7, 0.60, [("総コミット16.5億→", MUT, False, 10.5), ("MOIC 12.5X", GOLD, True, 11),
     ("／対ファンド ", MUT, False, 10.5), ("1.37X", CY, True, 11)], line_spacing=1.1)
# distribution rows right
rows = [("Solafune", "★唯一のFR（突出1本）", "206", GOLD),
        ("Humanity Brain", "勝ち馬B：実在候補8X（→論破③）", "136", CY),
        ("C", "中堅勝ち4.9X → ここで崖", "84", MUT),
        ("D–J", "2X以下に急減・大半0-1Xの薄い裾7社", "98", DIM)]
yy = 2.18
for name, tier, rec, c in rows:
    card(s, 3.85, yy, 5.55, 0.50)
    rect(s, 3.85, yy, 0.07, 0.50, c)
    text(s, 4.05, yy, 3.0, 0.50, name, size=12, color=WH, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 5.95, yy, 2.45, 0.50, tier, size=9.3, color=MUT, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 8.4, yy, 0.95, 0.50, f"¥{rec}億", size=13.5, color=c, bold=True, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
    yy += 0.56
# bottom math
trio = [("¥123億", "投下資本（10社）", CY), ("¥524億 / 4.3X", "グロス回収", GOLD), ("3.0X", "ネットDPI（検算一致）", PUR)]
x = 0.6
for b, l, c in trio:
    card(s, x, 4.46, 2.88, 0.70)
    text(s, x+0.18, 4.50, 2.5, 0.40, b, size=17, color=c, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x+0.18, 4.88, 2.5, 0.24, l, size=9.5, color=MUT)
    x += 2.93
footer(s, 12, "出典：02「ポートフォリオ」(net DPI≈3.0X) ／ →詳細：AP3（10社分布の崖・各社MOIC・感度4ケース）")

# =====================================================================
# S13 — 論破①「1.37XでFRは盛りすぎ？」（HP3→2）
# =====================================================================
rebuttal(13, 2, "論破① ・ DEBUNK ROUND 2（誇大）",
    "「1.37Xで“ファンドリターナー”は盛りすぎでは？」",
    "1.37Xって全然3.0Xに足りてないですよね？",
    [("別レイヤー＋鎖で返す。", WH, True), ("1.37Xは対ファンド比・MOIC12.5X、S12の逆算に飛躍ゼロ。", GOLD, True), ("どこが飛躍か挙げて。", WH, True)],
    [[("● ファンドリターナー＝1社で元手150億をほぼ1本ぶん回収＝206億／150億＝1.37X（対ファンド比）", MUT, False, 11.5)],
     [("● 投じた額で見れば：206億 ÷ 総コミット16.5億 ＝ ", MUT, False, 11.5), ("MOIC 12.5X", GOLD, True, 11.5), ("／全体3.0Xは10社分布で別途", MUT, False, 11.5)],
     [("● “全体が盛り”なら、S12で閉じた逆算の鎖のどこかが切れる＝DPI3.0X←206億←1,717億←17倍←101億は飛躍ゼロ", MUT, False, 11.5)],
     [("● だから問い返す：", DIM, False, 11), ("「どこが飛躍か、1つ挙げてください」", WH, True, 11.5), ("──挙がらなければ、盛りではなく定義です", DIM, False, 11)]],
    src="→詳細：AP3（対ファンド1.37X／対投資額MOIC12.5X／全体ネット3.0Xの3層＋逆算カスケード）")

# =====================================================================
# S14 — 論破③「Humanity Brain＝世論操作？」（HP2→1）
# =====================================================================
rebuttal(14, 3, "論破③ ・ DEBUNK ROUND 3（倫理）",
    "「Humanity Brainって“世論操作”の会社では？」",
    "認知戦AIって、世論を操作する危ない会社じゃないですか？",
    [("逆です。", WH, True), ("用途は“認知防衛（誤情報耐性）”に限定、", GOLD, True), ("攻撃的な利用は投資契約で禁止。", WH, True)],
    [[("● Humanity Brain＝認知戦対策AI・世論分布予測（防衛装備庁「社会シミュレーション」採択）", MUT, False, 11.5)],
     [("● 線引き：対象国の認知“操作”ではなく、自国の認知“防衛”＝誤情報への耐性に用途限定", MUT, False, 11.5)],
     [("● 投資契約に influence-ops の攻撃的利用禁止コベナンツ＋第三国移転制限", MUT, False, 11.5)],
     [("● Solafune（物理GEOINT）の隣＝認知ドメインの解析層。守る側です", DIM, False, 11)]],
    src="※前回唯一刺さったESGを塞ぐ ／ →詳細：AP7（用途限定・契約コベナンツ）")

# =====================================================================
# S15 — リスクと反論への備え
# =====================================================================
s = slide()
header(s, "RISK ・ リスクと備え", "残る弱点を、契約と設計で潰す")
keymsg(s, [("政府依存・チケット不整合・出口持続性に、", WH, True),
           ("契約・設計レベルの具体策を当てる", GOLD, True), ("（倫理は論破③で対応済）", MUT, True)])
risks = [("政府依存・調達変動", "マイルストン連動トランシェ＋同志国・民間へ需要分散"),
         ("実シード2億 vs 想定5.5億", "シード〜プレA単独リード拡大で再設計"),
         ("出口倍率の持続性", "Exit後ろ倒し＋国内M&A併行で備える"),
         ("倫理・ESG（SIGINT/LAWS）", "→ 論破③で用途限定＋コベナンツとして対応済")]
pos = [(0.6, 2.18), (5.1, 2.18), (0.6, 3.62), (5.1, 3.62)]
for (px, py), (t, d) in zip(pos, risks):
    card(s, px, py, 4.30, 1.34)
    text(s, px+0.24, py+0.16, 3.9, 0.36, [("⚠ ", GOLD, True, 13), (t, WH, True, 13.5)])
    text(s, px+0.24, py+0.62, 3.85, 0.62, d, size=11, color=MUT, line_spacing=1.12)
footer(s, 15, "→詳細：AP8（リスク×緩和の全表・既知の弱点a〜e）")

# =====================================================================
# S16 — 決着（論破④なぜ今あなた達＝総回収・投資判断 / HP1→0）
# =====================================================================
set_theme("dark")
s = slide(BG_COVER)
text(s, 0.6, 0.42, 9.0, 0.30, "FINISH ・ 決着（論破④：なぜ今・あなた達）", size=12, color=RED, bold=True)
text(s, 0.6, 0.86, 9.0, 0.66, "論破ポイント、全部塞ぎました", size=33, color=WH, bold=True)
keymsg(s, [("なぜ私たち × なぜ今 × なぜSolafune ", WH, True),
           ("― 3つ、全部揃ってます。", GOLD, True)], y=1.66, h=0.58, barc=RED)
# HP cleared
text(s, 0.6, 2.44, 3.0, 0.30, [("BOSS HP ", DIM, True, 10), ("0/4 ", CY, True, 13), ("撃破", CY, True, 12)])
px = 2.6
for i in range(4):
    rect(s, px, 2.48, 0.34, 0.18, "243154"); px += 0.44
text(s, 4.9, 2.42, 4.5, 0.30, "誇大→定義／感想→統計／倫理→契約／なぜ今→3点回収", size=10.5, color=MUT)
# decision card
card(s, 0.6, 2.90, 8.8, 1.34, fill=KMFILL, line=GOLD, radius=0.10, shadow=False)
text(s, 0.86, 3.02, 8.3, 0.40, [("投資判断：", GOLD, True, 13), ("初回5.5億・持分約15%でリード、リザーブ厚めでフォローオン＝ファンドリターナー候補", WH, True, 12.5)])
text(s, 0.86, 3.46, 8.3, 0.40, [("ポジション約206億／対ファンド ", MUT, False, 11.5), ("1.37X", CY, True, 13),
     ("／ネットDPI ", MUT, False, 11.5), ("3.0X", GOLD, True, 13), ("／Exit ", MUT, False, 11.5), ("1,717億", WH, True, 12.5)])
text(s, 0.86, 3.88, 8.3, 0.30, [("次アクション：", MUT, False, 11), ("シード〜プレA単独リード組成の条件交渉、KPI連動トランシェ設計", MUT, False, 11)])
text(s, 0.6, 4.50, 9.0, 0.40, [("「……で、論破できました？」　", RED, True, 13), ("地球を読む者が、次の安全保障を制する。", CY, True, 14)])
text(s, 0.6, 5.06, 9.0, 0.28, "ご清聴ありがとうございました", size=10.5, color=DIM)
footer(s, 16)

# =====================================================================
# Appendix（AP1〜9・補足/QA用・本編では飛ばす・各0秒）／白テーマ
#   本編ボディ末尾の「→詳細：APx」から開く。フル台帳（数値・導出・出典・仮印）はここ。
# =====================================================================
set_theme("light")

appendix(1, "comps・EV/Sales 17倍の導出", [
    ("H", "採用comps群（バブル除外・フォワード基準）"),
    "Palantirプレ局面20倍／Anduril14倍（$61B÷2026予想$4.3B）／BlackSky10倍／Synspective7倍／Planet6倍",
    [("→ ", CY, True, 9.5), ("中央値10倍・平均11.4倍・レンジ6〜20倍", WH, True, 9.5)],
    ("H", "採用17倍の導出（中央値10→17＝+7）と除外"),
    "要因分解＝ソフト比率+3／粗利優位+2／主権的信認+2（主権だけで1.7倍ではない）",
    "バブル除外：Palantir現況60.7倍（歴史的中央値25倍を143%超）・Planet現況29倍（宇宙ラリー）",
    "剥落リスク：トレーリングPSR7.1倍→フォワード約2.7倍（Synspective例）／持続性は感度12〜20倍で吸収",
    ("H", "comps評価額一覧"),
    "Palantir約3,060〜3,080億$・60.7倍／Anduril$61B／Helsing$18B／Shield AI$12.7B",
    "国内：Synspective約1,700億（1社1,500億Exit実証）／iQPS約1,216億／Axelspace2025/8上場／防衛テックVC2025＝491億$",
], src="出典：01前提パック §4・§5（S8・S10・S11の裏付け）")

appendix(2, "市場ボトムアップ（TAM／SAM／SOM内訳）", [
    ("H", "TAM（レイヤー別）"),
    "GeoAI 371→629億$（CAGR11.1%）／防衛地理空間1,503→2,821億$／地理空間アナリティクス1,143→2,265億$",
    ("H", "SAM＝年約170億（公的予算×解析ソフト充当率8〜15%［想定］）"),
    "PFI(405×8%)≈32＋宇宙基金(200×15%)≈30＋防衛AI(500×10%)≈50＋経済安保(300×10%)≈30＋その他官庁(200×15%)≈30",
    ("H", "SOM 5セグメント（誰に・いくらで・何件｜2030→2032）"),
    "A 日本政府 2.5億/件 → 12件30億 → 14件35億",
    "B 宇宙基金 3億/件 → 5件15億 → 6件18億　／　C 同志国MoU 1.5億/件 → 4件6億 → 6件9億（2026/4輸出緩和が追い風）",
    "D 経済安保 2億/件 → 7件14億 → 9件18億　／　E 民間 2億/件 → 10件20億 → 11件21億",
    [("計＝", MUT, False, 9.5), ("38件85億（2030）→ 46件101億（2032）", GOLD, True, 9.5),
     ("。確度 A＞B/D＞C（有償実績ゼロ＝最優先で実証）／感度 81〜121億（中心101）", MUT, False, 9.5)],
], src="出典：01前提パック §2・§6（S7・S9の裏付け）")

appendix(3, "リターン分布（power-law 10社）＆感度", [
    ("H", "10社のべき乗則分布（リザーブは勝ち馬に集中・下位3社は初回チケットのみ）"),
    "Solafune206億(12.5X)＋HB136億(8X)＋中堅C84億(4.9X) →【明確な崖】→ 44億(2.8X)＋26億(1.7X)＋16億(1.2X)＋0-1X〜全損4社",
    [("→ 投下約123億 → グロス回収約524億 → グロス4.3X → ", MUT, False, 9.5),
     ("ネットDPI約3.0X", GOLD, True, 9.5), ("（上位2社で回収の約65%）", MUT, False, 9.5)],
    ("H", "3つの倍率（混同が誤読の正体）"),
    "対ファンド1.37X（206/150）／対投資額MOIC12.5X（206/16.5）／全体ネット3.0X",
    ("H", "感度4ケース（保守域 ≥15倍 & ≥10% で安定して○）"),
    "①17×12%=206億○(+37%)　②15×11%=167億○(+11%)　③15×10%=151億○(+1%)　④12×8%=97億×(−35%)",
    "×は④同時悪化のみ＝Exit後ろ倒し(2034)で吸収／勝ち馬B＝Humanity Brain＝防衛装備庁「社会シミュレーション」採択［要検証］",
], src="出典：02財務モデル（ポートフォリオ/リターン）・01 §1（S12・S13の裏付け）")

appendix(4, "ファンド設計詳細（LP構成・件数算術）", [
    ("H", "LP構成 40/30/30 の合意ロジック［設計値］"),
    "事業会社40%（事業理解）→集中13%を許容／年金・大学30%（超長期）→運用10年+2・Exit後ろ倒しを許容",
    "金融機関30%（規律）→13%を上限にキャップ",
    [("→ ", CY, True, 9.5), ("集中度・期間はGP願望でなくLP合意済みの制約", WH, True, 9.5)],
    ("H", "件数算術・Pre逆算"),
    "GP1人が深く関与できるのは5社未満＝GP2で約10社／投資可能125億÷(チケット5.5＋リザーブ厚め)≒10〜11社",
    "Pre31＋初回5.5＝Post36.5、5.5/36.5＝15.1%（pre/postは非開示ゆえファンド想定からの逆算）",
], src="出典：01前提パック §0・§0-1・§8（S4の裏付け）")

appendix(5, "Solafune moat反証＆受注実績", [
    ("H", "moat4点＋反証条件"),
    "①センサー非依存のマルチソース解析（反証：ハード内製化／Palantir日本参入）",
    "②120カ国コミュニティ（反証：収益化せず）／③政府・国際機関リレーション（反証：スポット止まり）",
    "④自然言語UI Shigenbot（反証：汎用LLM+GISで模倣可能になれば薄れる）",
    ("H", "受注・MoU実績"),
    "防衛省（2024/9＋2026/1 AI2件）・内閣府・警察庁受注",
    "UNIDO $5.34M（コンゴ民SGN-C MoU）・セネガル/エジプト/ガーナ展開／UNDP掲載のSHIGEN AI・ShigenBot",
], src="出典：01前提パック §3・§8、inputs/Solafune調査.md（S6の裏付け）")

appendix(6, "公的需要プール（金額）と制度改革", [
    ("H", "公的需要プール（金額）"),
    "防衛費2025年度8.7兆円（2026案8.8兆＝再編込み9兆353億）／衛星コンステPFI約2,831億（2026-2031・7社落札）",
    "宇宙戦略基金10年総額1兆円／経済安保SC強靱化（令和4＝1兆358億・令和5＝9,172億・令和6＝2,300億）",
    "日米重要鉱物アクションプラン（2026/3発出）／安保技術研究推進制度 3億(2015)→約110億",
    ("H", "制度改革3点（2026）"),
    "①輸出緩和（2026/4/21・協定国へ移転可・戦闘中の国は除外）②ファストパス調達（約3.5か月）③前払い・部分払い",
], src="出典：01前提パック §6・§8（S5の裏付け）")

appendix(7, "ESG・契約コベナンツ詳細", [
    ("H", "用途限定（Humanity Brain）"),
    "認知「防衛」（誤情報耐性）に限定＝対象国の認知「操作」ではない",
    "位置づけ＝Solafune（物理GEOINT）隣の認知ドメイン解析層＝べき乗則の勝ち馬B（7-8X級）",
    ("H", "投資契約コベナンツ"),
    "influence-ops攻撃的利用禁止＋第三国移転制限＋LAWS非関与・輸出管理遵守・データ主権",
    "年金・大学基金LPのガバナンス要件とも整合",
], src="出典：01前提パック §3・§7・§0-1（S14・S15の裏付け）")

appendix(8, "リスク×緩和の全表＆既知の弱点", [
    ("H", "リスク×緩和策（全リスト）"),
    "①政府依存・調達サイクル → マイルストン連動トランシェ＋同志国・民間へ需要分散",
    "②チケット不整合（実シード2億 vs 想定5.5億）→ シード〜プレA単独リード拡大で再設計（post36.5/15%）",
    "③出口不確実性（バブル持続性＝Rheinmetall CEO警告）→ Exit後ろ倒し(2034)＋国内プライム/商社M&A併行",
    "④競争・人材 → マルチソース解析＋主権的信認＋120カ国コミュニティの採用ファネル",
    ("H", "既知の弱点（QA頻出）"),
    "a 実シード2億vs5.5億／b 倍率持続性／c Post-Val交渉の蓋然性／d 政府依存／e 1社集中",
], src="出典：01前提パック §7、STATUS §5（S15の裏付け）")

appendix(9, "チーム・付加価値・deal access詳細", [
    ("H", "チーム（「宇宙×防衛×資源」の当事者）"),
    "宇宙VC1号（裾野・部品/データ）2.0X／2号（宇宙データ・解析）1.5X",
    "安保クリアランス保持・防衛装備庁・国際機関（UNIDO/UNDP）・資源商社OBのオペレーティングパートナー（設計値）",
    ("H", "付加価値4点（勝たせる）／deal access4点（入れてもらえる）"),
    "付加価値：①取締役就任②120カ国コミュニティ接続③防衛装備庁・国際機関で政府アンカー受注④リザーブ55%フォローオン",
    "deal access：①クリアランス取締役＋防衛装備庁で政府アンカー後押し②国産GP＝主権調達適格（外資不可）③商社OBで販路・M&A買い手④フォローオンで資本確実性",
    [("load-bearing：", RED, True, 9.5), ("崩れるとS4持分15%・S12の206億/1.37Xが崩れる＝リターン論の前提", WH, True, 9.5)],
], src="出典：01前提パック §0・§1（S1・S2・S3の裏付け）")

prs.save(OUT)
print("saved:", OUT, "slides:", len(prs.slides._sldIdLst))
